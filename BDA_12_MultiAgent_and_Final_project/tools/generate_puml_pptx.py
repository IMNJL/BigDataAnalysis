#!/usr/bin/env python3
"""Render PlantUML file to PNG via plantuml server and create a PPTX embedding it.

Usage:
    python tools/generate_puml_pptx.py

Outputs:
 - presentation/architecture.png
 - presentation/experiment_presentation.pptx
"""
import os
import sys
import zlib
import base64
import requests

PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
PUML_PATH = os.path.join(PROJECT_ROOT, 'diagrams', 'system_architecture.puml')
OUT_IMG = os.path.join(PROJECT_ROOT, 'presentation', 'architecture.png')
OUT_PPTX = os.path.join(PROJECT_ROOT, 'presentation', 'experiment_presentation.pptx')

PLANTUML_SERVER = os.environ.get('PLANTUML_SERVER', 'https://www.plantuml.com/plantuml/png')


# PlantUML requires a special base64 encoding after zlib(deflate).
# Implementation based on PlantUML server encoding.
ENC_MAP = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz-_"


def encode_plantuml(text: str) -> str:
    data = text.encode('utf-8')
    compressed = zlib.compress(data, 9)[2:-4]
    res = ''
    b = compressed
    i = 0
    while i < len(b):
        if i+3 <= len(b):
            v = (b[i] << 16) | (b[i+1] << 8) | b[i+2]
            res += ENC_MAP[(v >> 18) & 0x3f]
            res += ENC_MAP[(v >> 12) & 0x3f]
            res += ENC_MAP[(v >> 6) & 0x3f]
            res += ENC_MAP[v & 0x3f]
            i += 3
        else:
            # handle remaining bytes
            rem = b[i:]
            v = 0
            for x in rem:
                v = (v << 8) | x
            # shift to make 3 bytes
            v <<= (3 - len(rem)) * 8
            res += ENC_MAP[(v >> 18) & 0x3f]
            res += ENC_MAP[(v >> 12) & 0x3f]
            if len(rem) == 1:
                res += ENC_MAP[(v >> 6) & 0x3f]
                res += ENC_MAP[v & 0x3f]
            elif len(rem) == 2:
                res += ENC_MAP[(v >> 6) & 0x3f]
                res += ENC_MAP[v & 0x3f]
            i = len(b)
    return res


def fetch_plantuml_png(puml_text: str, out_path: str) -> bool:
    code = encode_plantuml(puml_text)
    url = f"{PLANTUML_SERVER}/{code}"
    print('Requesting PlantUML server URL:', url)
    r = requests.get(url, stream=True)
    if r.status_code != 200:
        print('PlantUML server returned', r.status_code)
        return False
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, 'wb') as fh:
        for chunk in r.iter_content(1024*16):
            fh.write(chunk)
    print('Saved diagram to', out_path)
    return True


def make_pptx_with_image(img_path: str, out_pptx: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        print('python-pptx not available:', e)
        print('Install with: python -m pip install python-pptx Pillow')
        return False

    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = 'Multi-Agent RAG Experiment'
    slide.placeholders[1].text = 'Architecture & Flow'

    # Image slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = top = Inches(0.5)
    max_width = Inches(9)
    slide.shapes.add_picture(img_path, left, top, width=max_width)

    os.makedirs(os.path.dirname(out_pptx), exist_ok=True)
    prs.save(out_pptx)
    print('Saved presentation to', out_pptx)
    return True


def main():
    if not os.path.exists(PUML_PATH):
        print('PUML source not found at', PUML_PATH)
        sys.exit(2)
    with open(PUML_PATH, 'r', encoding='utf-8') as fh:
        puml = fh.read()
    ok = fetch_plantuml_png(puml, OUT_IMG)
    if not ok:
        print('Failed to render PlantUML via server. You can render locally with PlantUML jar.')
    else:
        make_pptx_with_image(OUT_IMG, OUT_PPTX)


if __name__ == '__main__':
    main()
