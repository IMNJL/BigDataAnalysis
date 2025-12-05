#!/usr/bin/env python3
"""Render all PlantUML files in `diagrams/` to PNG via PlantUML server and create a PPTX embedding them.

Usage:
    python tools/generate_all_puml_pptx.py

Outputs:
 - presentation/<puml_basename>.png for each .puml
 - presentation/architecture_all_slides.pptx
"""
import os
import sys
import zlib
import requests
from glob import glob

PROJECT_ROOT = os.path.dirname(os.path.dirname(__file__))
DIAGRAMS_DIR = os.path.join(PROJECT_ROOT, 'diagrams')
OUT_DIR = os.path.join(PROJECT_ROOT, 'presentation')
PLANTUML_SERVER = os.environ.get('PLANTUML_SERVER', 'https://www.plantuml.com/plantuml/png')

ENC_MAP = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz-_"


def encode_plantuml(text: str) -> str:
    data = text.encode('utf-8')
    compressed = zlib.compress(data, 9)[2:-4]
    res = ''
    b = compressed
    i = 0
    while i < len(b):
        if i+3 <= len(b):
            v = (b[i] << 16) | (b[i+1] << 8) | b[i+2]
            res += ENC_MAP[(v >> 18) & 0x3f]
            res += ENC_MAP[(v >> 12) & 0x3f]
            res += ENC_MAP[(v >> 6) & 0x3f]
            res += ENC_MAP[v & 0x3f]
            i += 3
        else:
            rem = b[i:]
            v = 0
            for x in rem:
                v = (v << 8) | x
            v <<= (3 - len(rem)) * 8
            res += ENC_MAP[(v >> 18) & 0x3f]
            res += ENC_MAP[(v >> 12) & 0x3f]
            if len(rem) == 1:
                res += ENC_MAP[(v >> 6) & 0x3f]
                res += ENC_MAP[v & 0x3f]
            elif len(rem) == 2:
                res += ENC_MAP[(v >> 6) & 0x3f]
                res += ENC_MAP[v & 0x3f]
            i = len(b)
    return res


def fetch_plantuml_png(puml_text: str, out_path: str) -> bool:
    code = encode_plantuml(puml_text)
    url = f"{PLANTUML_SERVER}/{code}"
    print('Requesting', url)
    r = requests.get(url, stream=True)
    if r.status_code != 200:
        print('PlantUML server returned', r.status_code)
        return False
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, 'wb') as fh:
        for chunk in r.iter_content(1024*16):
            fh.write(chunk)
    print('Saved diagram to', out_path)
    return True


def make_pptx(image_paths, out_pptx):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        print('python-pptx not available:', e)
        print('Install with: python -m pip install python-pptx Pillow')
        return False

    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = 'Multi-Agent RAG System — Diagrams'
    slide.placeholders[1].text = 'Generated diagrams'

    for img in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = Inches(0.5)
        max_w = Inches(9)
        slide.shapes.add_picture(img, left, top, width=max_w)

    os.makedirs(os.path.dirname(out_pptx), exist_ok=True)
    prs.save(out_pptx)
    print('Saved presentation to', out_pptx)
    return True


def main():
    puml_files = glob(os.path.join(DIAGRAMS_DIR, '*.puml'))
    if not puml_files:
        print('No .puml files in', DIAGRAMS_DIR)
        sys.exit(2)

    images = []
    for p in puml_files:
        name = os.path.splitext(os.path.basename(p))[0]
        out_img = os.path.join(OUT_DIR, f"{name}.png")
        with open(p, 'r', encoding='utf-8') as fh:
            text = fh.read()
        ok = fetch_plantuml_png(text, out_img)
        if ok:
            images.append(out_img)

    if images:
        make_pptx(images, os.path.join(OUT_DIR, 'architecture_all_slides.pptx'))


if __name__ == '__main__':
    main()
